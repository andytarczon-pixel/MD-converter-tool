import base64
import anthropic
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None
    MSO_SHAPE_TYPE = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

HAIKU_MODEL = "claude-haiku-4-5-20251001"
INPUT_COST_PER_MTOK = 0.80   # $0.80 per million input tokens
OUTPUT_COST_PER_MTOK = 4.00  # $4.00 per million output tokens
EST_INPUT_TOKENS_PER_IMAGE = 1500
EST_OUTPUT_TOKENS_PER_IMAGE = 200
MIN_IMAGE_BYTES = 2000  # skip tiny/decorative images


class FileConverter:
    def __init__(self, api_key, images_enabled=True, progress_callback=None, log_callback=None):
        self.images_enabled = images_enabled
        self.client = anthropic.Anthropic(api_key=api_key) if (images_enabled and api_key) else None
        self.progress_callback = progress_callback
        self.log_callback = log_callback

    def log(self, msg):
        if self.log_callback:
            self.log_callback(msg)

    def progress(self, value):
        if self.progress_callback:
            self.progress_callback(min(max(value, 0.0), 1.0))

    # ── Cost estimation ──────────────────────────────────────────────────────

    def estimate_images(self, file_paths):
        total = 0
        for fp in file_paths:
            ext = Path(fp).suffix.lower()
            if ext in ('.pptx', '.ppt'):
                total += self._count_pptx_images(fp)
            elif ext == '.pdf':
                total += self._count_pdf_images(fp)
        return total

    def estimate_cost(self, image_count):
        input_tok = image_count * EST_INPUT_TOKENS_PER_IMAGE
        output_tok = image_count * EST_OUTPUT_TOKENS_PER_IMAGE
        return (input_tok * INPUT_COST_PER_MTOK / 1_000_000 +
                output_tok * OUTPUT_COST_PER_MTOK / 1_000_000)

    def _count_pptx_images(self, path):
        if not Presentation:
            return 0
        try:
            prs = Presentation(path)
            return sum(
                1 for slide in prs.slides
                for shape in slide.shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )
        except Exception:
            return 0

    def _count_pdf_images(self, path):
        if not fitz:
            return 0
        try:
            doc = fitz.open(path)
            count = sum(len(page.get_images()) for page in doc)
            doc.close()
            return count
        except Exception:
            return 0

    # ── Main entry point ─────────────────────────────────────────────────────

    def convert_file(self, file_path, output_folder, overwrite_mode="overwrite"):
        ext = Path(file_path).suffix.lower()
        stem = Path(file_path).stem
        output_path = Path(output_folder) / f"{stem}.md"

        if output_path.exists() and overwrite_mode == "skip":
            return False, str(output_path), "Skipped (file exists)"

        if ext in ('.pptx', '.ppt'):
            if not Presentation:
                return False, None, "python-pptx not installed"
            md = self._convert_pptx(file_path)
        elif ext == '.pdf':
            if not fitz:
                return False, None, "PyMuPDF not installed"
            md = self._convert_pdf(file_path)
        else:
            return False, None, f"Unsupported format: {ext}"

        output_path.write_text(md, encoding='utf-8')
        return True, str(output_path), "Success"

    # ── Image description via Claude Haiku ───────────────────────────────────

    def _detect_media_type(self, data):
        if data[:8] == b'\x89PNG\r\n\x1a\n':
            return "image/png"
        if data[:2] == b'\xff\xd8':
            return "image/jpeg"
        if data[:6] in (b'GIF87a', b'GIF89a'):
            return "image/gif"
        if data[:4] == b'RIFF' and data[8:12] == b'WEBP':
            return "image/webp"
        return "image/png"

    def _describe_image(self, image_bytes, context=""):
        if not self.images_enabled or not self.client:
            return "[Image — description disabled]"
        if not image_bytes or len(image_bytes) < MIN_IMAGE_BYTES:
            return None

        media_type = self._detect_media_type(image_bytes)
        b64 = base64.standard_b64encode(image_bytes).decode('utf-8')

        try:
            response = self.client.messages.create(
                model=HAIKU_MODEL,
                max_tokens=400,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": b64
                            }
                        },
                        {
                            "type": "text",
                            "text": (
                                f"Context: {context}\n\n"
                                "Describe this image concisely. "
                                "Extract ALL visible text verbatim. "
                                "For charts, graphs, or tables, describe the data shown. "
                                "For diagrams, describe the structure and labels."
                            )
                        }
                    ]
                }]
            )
            return response.content[0].text.strip()
        except Exception as e:
            return f"[Image description unavailable: {e}]"

    # ── PPTX conversion ──────────────────────────────────────────────────────

    def _convert_pptx(self, file_path):
        prs = Presentation(file_path)
        lines = []

        title = Path(file_path).stem
        try:
            if prs.core_properties.title and prs.core_properties.title.strip():
                title = prs.core_properties.title.strip()
        except Exception:
            pass

        lines.append(f"# {title}\n")
        lines.append(f"*Source: {Path(file_path).name}*\n")
        lines.append("---\n")

        total = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, 1):
            self.log(f"Slide {slide_num}/{total}...")
            self.progress((slide_num - 1) / total)

            slide_title = ""
            title_shape = None
            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    slide_title = slide.shapes.title.text.strip()
                    title_shape = slide.shapes.title
            except Exception:
                pass

            heading = f"## Slide {slide_num}"
            if slide_title:
                heading += f": {slide_title}"
            lines.append(f"\n{heading}\n")

            # Speaker notes
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"\n> *Speaker Notes: {notes}*\n")
            except Exception:
                pass

            for shape in slide.shapes:
                if shape == title_shape:
                    continue

                if shape.has_text_frame:
                    self._append_pptx_text(shape.text_frame, lines)

                elif shape.has_table:
                    self._append_pptx_table(shape.table, lines)

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    self.log(f"Slide {slide_num}/{total}: describing image...")
                    desc = self._describe_image(
                        shape.image.blob,
                        context=f"Slide {slide_num} of '{title}'"
                    )
                    if desc:
                        lines.append(f"\n> **[Image]** {desc}\n")

                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for s in shape.shapes:
                        if s.has_text_frame:
                            self._append_pptx_text(s.text_frame, lines)
                        elif s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            self.log(f"Slide {slide_num}/{total}: describing grouped image...")
                            desc = self._describe_image(
                                s.image.blob,
                                context=f"Slide {slide_num} of '{title}'"
                            )
                            if desc:
                                lines.append(f"\n> **[Image]** {desc}\n")

        self.progress(1.0)
        return "\n".join(lines)

    def _append_pptx_text(self, text_frame, lines):
        for para in text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            level = para.level
            if level == 0:
                lines.append(f"\n{text}")
            else:
                indent = "  " * (level - 1)
                lines.append(f"{indent}- {text}")

    def _append_pptx_table(self, table, lines):
        lines.append("")
        if not table.rows:
            return
        headers = [cell.text.strip() for cell in table.rows[0].cells]
        lines.append("| " + " | ".join(headers) + " |")
        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in list(table.rows)[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")

    # ── PDF conversion ───────────────────────────────────────────────────────

    def _convert_pdf(self, file_path):
        doc = fitz.open(file_path)
        lines = []

        title = Path(file_path).stem
        try:
            meta = doc.metadata
            if meta.get('title') and meta['title'].strip():
                title = meta['title'].strip()
        except Exception:
            pass

        lines.append(f"# {title}\n")
        lines.append(f"*Source: {Path(file_path).name}*\n")
        lines.append("---\n")

        total = len(doc)

        for page_num, page in enumerate(doc, 1):
            self.log(f"Page {page_num}/{total}...")
            self.progress((page_num - 1) / total)

            lines.append(f"\n## Page {page_num}\n")

            # Extract structured text
            blocks = page.get_text("dict", sort=True)["blocks"]
            seen_xrefs = set()

            for block in blocks:
                btype = block.get("type")

                if btype == 0:  # text block
                    block_lines = []
                    for line in block.get("lines", []):
                        text = " ".join(s["text"] for s in line.get("spans", [])).strip()
                        if text:
                            block_lines.append(text)
                    if block_lines:
                        lines.append("\n".join(block_lines))

                elif btype == 1:  # inline image block
                    xref = block.get("image")
                    if xref and xref not in seen_xrefs:
                        seen_xrefs.add(xref)
                        try:
                            img_data = doc.extract_image(xref)
                            image_bytes = img_data["image"]
                            if len(image_bytes) >= MIN_IMAGE_BYTES:
                                self.log(f"Page {page_num}/{total}: describing image...")
                                desc = self._describe_image(
                                    image_bytes,
                                    context=f"Page {page_num} of '{title}'"
                                )
                                if desc:
                                    lines.append(f"\n> **[Image]** {desc}\n")
                        except Exception:
                            pass

            # Catch any images not surfaced in blocks
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                try:
                    img_data = doc.extract_image(xref)
                    image_bytes = img_data["image"]
                    if len(image_bytes) >= MIN_IMAGE_BYTES:
                        self.log(f"Page {page_num}/{total}: describing image...")
                        desc = self._describe_image(
                            image_bytes,
                            context=f"Page {page_num} of '{title}'"
                        )
                        if desc:
                            lines.append(f"\n> **[Image]** {desc}\n")
                except Exception:
                    pass

        doc.close()
        self.progress(1.0)
        return "\n".join(lines)
